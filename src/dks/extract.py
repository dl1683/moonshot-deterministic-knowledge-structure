"""Claim extraction from unstructured text.

This module provides the Extractor protocol and default implementations for
extracting structured ClaimCore instances from text. The commitment boundary
is explicit: ExtractionResult is non-deterministic output. Only when claims
are committed via store.assert_revision() do they become deterministic data.

Includes:
- RegexExtractor: Zero-dependency pattern-based extraction
- LLMExtractor: LLM-backed open-domain extraction
- TextChunker: Smart text splitting with overlap for document ingestion
- PDFExtractor: PDF text extraction + chunking via PyMuPDF
- DocxExtractor: Word (.docx) extraction via python-docx
- PptxExtractor: PowerPoint (.pptx) extraction via python-pptx
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Protocol, runtime_checkable

from .core import ClaimCore, Provenance, canonicalize_text


@dataclass(frozen=True)
class ExtractionResult:
    """Result of extracting claims from text.

    This is the output of the non-deterministic extraction phase.
    Claims become deterministic only after commitment to a KnowledgeStore.
    """
    claims: tuple[ClaimCore, ...]
    provenance: tuple[Provenance, ...]
    raw_text: str
    metadata: dict[str, Any] = field(default_factory=dict)


@runtime_checkable
class Extractor(Protocol):
    """Protocol for claim extraction backends.

    Implementations may be deterministic (regex) or non-deterministic (LLM).
    """

    def extract(
        self,
        text: str,
        *,
        claim_types: list[str] | None = None,
    ) -> ExtractionResult:
        """Extract structured claims from text.

        Args:
            text: Unstructured input text.
            claim_types: Optional filter — only extract claims of these types.

        Returns:
            ExtractionResult with extracted claims and provenance.
        """
        ...


class RegexExtractor:
    """Zero-dependency regex-based extractor for structured patterns.

    Extracts claims matching registered patterns. Good for:
    - Dates, numbers, known templates
    - Structured log lines
    - Key-value pairs

    Not suitable for open-domain extraction (use LLMExtractor for that).
    """

    def __init__(self) -> None:
        self._patterns: list[tuple[str, re.Pattern[str], list[str]]] = []

    def register_pattern(
        self,
        claim_type: str,
        pattern: str,
        slot_names: list[str],
    ) -> None:
        """Register a regex pattern for a claim type.

        Named groups in the pattern are mapped to slot_names in order.
        If the pattern uses unnamed groups, they're mapped positionally.
        """
        self._patterns.append((claim_type, re.compile(pattern), slot_names))

    def extract(
        self,
        text: str,
        *,
        claim_types: list[str] | None = None,
    ) -> ExtractionResult:
        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for claim_type, pattern, slot_names in self._patterns:
            if claim_types is not None and claim_type not in claim_types:
                continue

            for match in pattern.finditer(text):
                groups = match.groupdict() if match.groupdict() else {}
                if not groups:
                    # Use positional groups
                    groups = {
                        name: val
                        for name, val in zip(slot_names, match.groups())
                        if val is not None
                    }
                else:
                    # Map named groups to slot_names
                    mapped = {}
                    for name in slot_names:
                        if name in groups and groups[name] is not None:
                            mapped[name] = groups[name]
                    groups = mapped

                if groups:
                    claim = ClaimCore(
                        claim_type=claim_type,
                        slots={k: canonicalize_text(v) for k, v in groups.items()},
                    )
                    claims.append(claim)
                    provenances.append(Provenance(
                        source=f"regex:{claim_type}",
                        evidence_ref=match.group(0),
                    ))

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=text,
        )


class LLMExtractor:
    """LLM-backed extractor for open-domain claim extraction.

    Requires an LLM callable that takes a prompt and returns structured JSON.
    The LLM is called outside the deterministic boundary — results are
    non-deterministic until committed to a KnowledgeStore.

    Default model recommendation: Qwen3-0.6B (from model registry)
    for fast iteration. Upgrade to Qwen3-4B for production quality.
    """

    def __init__(
        self,
        llm_fn: Any,
        *,
        system_prompt: str | None = None,
        model_id: str = "Qwen/Qwen3-0.6B",
    ) -> None:
        """
        Args:
            llm_fn: Callable that takes (prompt: str) -> str (JSON response).
            system_prompt: Optional system prompt for the LLM.
            model_id: Model identifier for provenance tracking.
        """
        self._llm_fn = llm_fn
        self._system_prompt = system_prompt or self._default_system_prompt()
        self._model_id = model_id

    @staticmethod
    def _default_system_prompt() -> str:
        return (
            "Extract factual claims from the following text. "
            "Return a JSON array where each element has:\n"
            '  {"claim_type": "<type>", "slots": {"<role>": "<value>", ...}}\n'
            "Only extract concrete, verifiable facts. "
            "Use lowercase for all values."
        )

    def extract(
        self,
        text: str,
        *,
        claim_types: list[str] | None = None,
    ) -> ExtractionResult:
        import json

        type_filter = ""
        if claim_types:
            type_filter = f"\nOnly extract these claim types: {', '.join(claim_types)}"

        prompt = f"{self._system_prompt}{type_filter}\n\nText:\n{text}"

        raw_response = self._llm_fn(prompt)

        try:
            parsed = json.loads(raw_response)
        except (json.JSONDecodeError, TypeError):
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text=text,
                metadata={"error": "failed to parse LLM response", "raw": str(raw_response)},
            )

        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        items = parsed if isinstance(parsed, list) else [parsed]
        for item in items:
            if not isinstance(item, dict):
                continue
            ct = item.get("claim_type", "")
            slots = item.get("slots", {})
            if not ct or not slots:
                continue
            if not isinstance(slots, dict):
                continue

            claim = ClaimCore(
                claim_type=canonicalize_text(ct),
                slots={canonicalize_text(k): canonicalize_text(str(v)) for k, v in slots.items()},
            )
            claims.append(claim)
            provenances.append(Provenance(
                source=f"llm:{self._model_id}",
                evidence_ref=text[:200],
            ))

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=text,
            metadata={"model_id": self._model_id},
        )


class TextChunker:
    """Split text into overlapping chunks for document ingestion.

    Uses paragraph boundaries when possible, falls back to sentence
    boundaries, then character-level splitting. Chunks maintain overlap
    for context continuity.
    """

    def __init__(
        self,
        chunk_size: int = 1000,
        overlap: int = 200,
        min_chunk: int = 50,
    ) -> None:
        self._chunk_size = chunk_size
        self._overlap = overlap
        self._min_chunk = min_chunk

    def chunk(self, text: str) -> list[str]:
        """Split text into overlapping chunks.

        Strategy:
        1. Split on double newlines (paragraph boundaries)
        2. Merge small paragraphs to reach target chunk_size
        3. Split oversized paragraphs on sentence boundaries
        4. Apply overlap between consecutive chunks
        """
        if not text or not text.strip():
            return []

        # Split into paragraphs
        paragraphs = re.split(r'\n\s*\n', text)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        if not paragraphs:
            return []

        # Split oversized paragraphs into sentences
        segments: list[str] = []
        for para in paragraphs:
            if len(para) <= self._chunk_size:
                segments.append(para)
            else:
                # Split on sentence boundaries
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    if sent.strip():
                        segments.append(sent.strip())

        # Merge segments into chunks of target size
        chunks: list[str] = []
        current: list[str] = []
        current_len = 0

        for seg in segments:
            seg_len = len(seg)

            if current_len + seg_len + 1 > self._chunk_size and current:
                chunk_text = '\n\n'.join(current)
                if len(chunk_text) >= self._min_chunk:
                    chunks.append(chunk_text)
                current = [seg]
                current_len = seg_len
            else:
                current.append(seg)
                current_len += seg_len + 1

        # Final chunk
        if current:
            chunk_text = '\n\n'.join(current)
            if len(chunk_text) >= self._min_chunk:
                chunks.append(chunk_text)

        # Apply overlap: prepend tail of previous chunk to next chunk
        if self._overlap > 0 and len(chunks) > 1:
            overlapped: list[str] = [chunks[0]]
            for i in range(1, len(chunks)):
                prev = chunks[i - 1]
                overlap_text = prev[-self._overlap:] if len(prev) > self._overlap else prev
                # Find a word boundary in the overlap
                space_idx = overlap_text.find(' ')
                if space_idx > 0:
                    overlap_text = overlap_text[space_idx + 1:]
                overlapped.append(overlap_text + '\n\n' + chunks[i])
            chunks = overlapped

        return chunks


class PDFExtractor:
    """Extract text from PDFs and chunk into claims for ingestion.

    Uses PyMuPDF (fitz) for PDF text extraction. Each chunk becomes a
    ClaimCore with claim_type="document.chunk@v1" containing:
    - source: filename
    - chunk_idx: sequential index
    - page_start: first page of chunk
    - text: the chunk text (canonicalized)

    Requires: pip install PyMuPDF
    """

    CLAIM_TYPE = "document.chunk@v1"

    def __init__(
        self,
        chunker: TextChunker | None = None,
        *,
        extract_metadata: bool = True,
    ) -> None:
        self._chunker = chunker or TextChunker()
        self._extract_metadata = extract_metadata

    def extract_pdf(self, path: str | Path) -> ExtractionResult:
        """Read a PDF file and extract chunked claims.

        Args:
            path: Path to the PDF file.

        Returns:
            ExtractionResult with one ClaimCore per chunk.
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF required: pip install PyMuPDF")

        path = Path(path)
        filename = path.name

        # Extract text from all pages
        doc = fitz.open(str(path))
        pages_text: list[tuple[int, str]] = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            if text and text.strip():
                pages_text.append((page_num, text))
        doc.close()

        if not pages_text:
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text="",
                metadata={"source": filename, "error": "no text extracted"},
            )

        # Combine all page text with page markers
        full_text = ""
        page_boundaries: list[tuple[int, int]] = []  # (char_offset, page_num)
        for page_num, text in pages_text:
            page_boundaries.append((len(full_text), page_num))
            full_text += text + "\n\n"

        # Clean up common PDF artifacts
        full_text = self._clean_pdf_text(full_text)

        # Chunk the text
        chunks = self._chunker.chunk(full_text)

        if not chunks:
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text=full_text[:500],
                metadata={"source": filename, "error": "no chunks produced"},
            )

        # Map chunks to page numbers
        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for idx, chunk_text in enumerate(chunks):
            # Find which page this chunk starts on
            chunk_start = full_text.find(chunk_text[:100])
            page_start = 0
            if chunk_start >= 0:
                for offset, pnum in page_boundaries:
                    if offset <= chunk_start:
                        page_start = pnum

            claim = ClaimCore(
                claim_type=self.CLAIM_TYPE,
                slots={
                    "source": canonicalize_text(filename),
                    "chunk_idx": str(idx),
                    "page_start": str(page_start),
                    "text": canonicalize_text(chunk_text[:200]),
                },
            )
            claims.append(claim)
            provenances.append(Provenance(
                source=f"pdf:{filename}",
                evidence_ref=chunk_text,
            ))

        metadata: dict[str, Any] = {
            "source": filename,
            "total_pages": len(pages_text),
            "total_chunks": len(chunks),
            "total_chars": len(full_text),
        }

        if self._extract_metadata and pages_text:
            try:
                doc = fitz.open(str(path))
                md = doc.metadata
                if md:
                    for key in ("title", "author", "subject"):
                        if md.get(key):
                            metadata[key] = md[key]
                doc.close()
            except (ImportError, ValueError, KeyError, RuntimeError, OSError):
                pass

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=full_text[:1000],
            metadata=metadata,
        )

    def extract(
        self,
        text: str,
        *,
        claim_types: list[str] | None = None,
    ) -> ExtractionResult:
        """Satisfy Extractor protocol — chunk text directly."""
        if claim_types and self.CLAIM_TYPE not in claim_types:
            return ExtractionResult(claims=(), provenance=(), raw_text=text)

        chunks = self._chunker.chunk(text)
        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for idx, chunk_text in enumerate(chunks):
            claim = ClaimCore(
                claim_type=self.CLAIM_TYPE,
                slots={
                    "source": "text_input",
                    "chunk_idx": str(idx),
                    "page_start": "0",
                    "text": canonicalize_text(chunk_text[:200]),
                },
            )
            claims.append(claim)
            provenances.append(Provenance(
                source="chunker",
                evidence_ref=chunk_text,
            ))

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=text[:1000],
            metadata={"total_chunks": len(chunks)},
        )

    @staticmethod
    def _clean_pdf_text(text: str) -> str:
        """Clean common PDF extraction artifacts."""
        # Remove excessive whitespace lines
        text = re.sub(r'\n{4,}', '\n\n\n', text)
        # Remove page headers/footers (lines with just numbers)
        text = re.sub(r'\n\s*\d+\s*\n', '\n', text)
        # Fix hyphenation at line breaks
        text = re.sub(r'(\w)-\n(\w)', r'\1\2', text)
        # Normalize whitespace within lines (but preserve paragraph breaks)
        lines = text.split('\n')
        cleaned = []
        for line in lines:
            line = ' '.join(line.split())
            cleaned.append(line)
        text = '\n'.join(cleaned)
        # Re-merge into paragraphs (single newlines become spaces within paragraphs)
        text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
        return text.strip()


class DocxExtractor:
    """Extract text from Word (.docx) documents and chunk into claims.

    Extracts paragraphs and table content, preserving document structure.
    Each chunk becomes a ClaimCore with claim_type="document.chunk@v1".

    Requires: pip install python-docx
    """

    CLAIM_TYPE = "document.chunk@v1"

    def __init__(
        self,
        chunker: TextChunker | None = None,
        *,
        extract_metadata: bool = True,
    ) -> None:
        self._chunker = chunker or TextChunker()
        self._extract_metadata = extract_metadata

    def extract_docx(self, path: str | Path) -> ExtractionResult:
        """Read a .docx file and extract chunked claims.

        Args:
            path: Path to the Word document.

        Returns:
            ExtractionResult with one ClaimCore per chunk.
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx required: pip install python-docx")

        path = Path(path)
        filename = path.name

        doc = Document(str(path))

        # Extract paragraph text
        sections: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                sections.append(text)

        # Extract table content (row by row)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    sections.append(row_text)

        if not sections:
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text="",
                metadata={"source": filename, "error": "no text extracted"},
            )

        full_text = "\n\n".join(sections)
        chunks = self._chunker.chunk(full_text)

        if not chunks:
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text=full_text[:500],
                metadata={"source": filename, "error": "no chunks produced"},
            )

        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for idx, chunk_text in enumerate(chunks):
            claim = ClaimCore(
                claim_type=self.CLAIM_TYPE,
                slots={
                    "source": canonicalize_text(filename),
                    "chunk_idx": str(idx),
                    "text": canonicalize_text(chunk_text[:200]),
                },
            )
            claims.append(claim)
            provenances.append(Provenance(
                source=f"docx:{filename}",
                evidence_ref=chunk_text,
            ))

        metadata: dict[str, Any] = {
            "source": filename,
            "total_paragraphs": len(doc.paragraphs),
            "total_tables": len(doc.tables),
            "total_chunks": len(chunks),
            "total_chars": len(full_text),
        }

        if self._extract_metadata:
            try:
                cp = doc.core_properties
                if cp.title:
                    metadata["title"] = cp.title
                if cp.author:
                    metadata["author"] = cp.author
                if cp.subject:
                    metadata["subject"] = cp.subject
            except (AttributeError, KeyError, ValueError):
                pass

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=full_text[:1000],
            metadata=metadata,
        )

    def extract(
        self,
        text: str,
        *,
        claim_types: list[str] | None = None,
    ) -> ExtractionResult:
        """Satisfy Extractor protocol — chunk text directly."""
        if claim_types and self.CLAIM_TYPE not in claim_types:
            return ExtractionResult(claims=(), provenance=(), raw_text=text)

        chunks = self._chunker.chunk(text)
        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for idx, chunk_text in enumerate(chunks):
            claim = ClaimCore(
                claim_type=self.CLAIM_TYPE,
                slots={
                    "source": "text_input",
                    "chunk_idx": str(idx),
                    "text": canonicalize_text(chunk_text[:200]),
                },
            )
            claims.append(claim)
            provenances.append(Provenance(
                source="chunker",
                evidence_ref=chunk_text,
            ))

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=text[:1000],
            metadata={"total_chunks": len(chunks)},
        )


class PptxExtractor:
    """Extract text from PowerPoint (.pptx) presentations and chunk into claims.

    Extracts text from slide shapes (text frames, tables, notes).
    Each chunk becomes a ClaimCore with claim_type="document.chunk@v1".

    Requires: pip install python-pptx
    """

    CLAIM_TYPE = "document.chunk@v1"

    def __init__(
        self,
        chunker: TextChunker | None = None,
        *,
        include_notes: bool = True,
    ) -> None:
        self._chunker = chunker or TextChunker()
        self._include_notes = include_notes

    def extract_pptx(self, path: str | Path) -> ExtractionResult:
        """Read a .pptx file and extract chunked claims.

        Args:
            path: Path to the PowerPoint file.

        Returns:
            ExtractionResult with one ClaimCore per chunk.
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required: pip install python-pptx")

        path = Path(path)
        filename = path.name

        prs = Presentation(str(path))
        slide_texts: list[tuple[int, str]] = []

        for slide_num, slide in enumerate(prs.slides):
            parts: list[str] = []

            # Extract text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)

                # Extract table content
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip(" |"):
                            parts.append(row_text)

            # Extract slide notes
            if self._include_notes and slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")

            if parts:
                slide_texts.append((slide_num, "\n".join(parts)))

        if not slide_texts:
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text="",
                metadata={"source": filename, "error": "no text extracted"},
            )

        # Combine slide text with slide markers
        full_text = ""
        slide_boundaries: list[tuple[int, int]] = []
        for slide_num, text in slide_texts:
            slide_boundaries.append((len(full_text), slide_num))
            full_text += text + "\n\n"
        full_text = full_text.strip()

        chunks = self._chunker.chunk(full_text)

        if not chunks:
            return ExtractionResult(
                claims=(),
                provenance=(),
                raw_text=full_text[:500],
                metadata={"source": filename, "error": "no chunks produced"},
            )

        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for idx, chunk_text in enumerate(chunks):
            # Find which slide this chunk starts on
            chunk_start = full_text.find(chunk_text[:100])
            slide_start = 0
            if chunk_start >= 0:
                for offset, snum in slide_boundaries:
                    if offset <= chunk_start:
                        slide_start = snum

            claim = ClaimCore(
                claim_type=self.CLAIM_TYPE,
                slots={
                    "source": canonicalize_text(filename),
                    "chunk_idx": str(idx),
                    "slide_start": str(slide_start),
                    "text": canonicalize_text(chunk_text[:200]),
                },
            )
            claims.append(claim)
            provenances.append(Provenance(
                source=f"pptx:{filename}",
                evidence_ref=chunk_text,
            ))

        metadata: dict[str, Any] = {
            "source": filename,
            "total_slides": len(prs.slides),
            "slides_with_text": len(slide_texts),
            "total_chunks": len(chunks),
            "total_chars": len(full_text),
        }

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=full_text[:1000],
            metadata=metadata,
        )

    def extract(
        self,
        text: str,
        *,
        claim_types: list[str] | None = None,
    ) -> ExtractionResult:
        """Satisfy Extractor protocol — chunk text directly."""
        if claim_types and self.CLAIM_TYPE not in claim_types:
            return ExtractionResult(claims=(), provenance=(), raw_text=text)

        chunks = self._chunker.chunk(text)
        claims: list[ClaimCore] = []
        provenances: list[Provenance] = []

        for idx, chunk_text in enumerate(chunks):
            claim = ClaimCore(
                claim_type=self.CLAIM_TYPE,
                slots={
                    "source": "text_input",
                    "chunk_idx": str(idx),
                    "text": canonicalize_text(chunk_text[:200]),
                },
            )
            claims.append(claim)
            provenances.append(Provenance(
                source="chunker",
                evidence_ref=chunk_text,
            ))

        return ExtractionResult(
            claims=tuple(claims),
            provenance=tuple(provenances),
            raw_text=text[:1000],
            metadata={"total_chunks": len(chunks)},
        )
